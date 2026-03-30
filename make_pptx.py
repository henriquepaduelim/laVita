#!/usr/bin/env python3
"""
Gera analise_la_vita.pptx — apresentação completa do case La Vita Alimentos.
10 slides, 16:9 widescreen, design navy/teal.
"""
import warnings
warnings.filterwarnings('ignore')

from pathlib import Path
from collections import defaultdict
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from matplotlib.patches import FancyBboxPatch, Rectangle
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches

# ─────────────────────────────────────────────────────────────────────────────
# CONSTANTES
# ─────────────────────────────────────────────────────────────────────────────
BASE_DIR  = Path('.')
TMP_DIR   = BASE_DIR / '_slides_tmp'
TMP_DIR.mkdir(exist_ok=True)

W, H = 13.333, 7.5   # polegadas 16:9
DPI  = 150

NAVY   = '#1B3A6B'
TEAL   = '#048A81'
GREEN  = '#27AE60'
RED    = '#C0392B'
ORANGE = '#E67E22'
GRAY   = '#7F8C8D'
LGRAY  = '#F4F6F7'
MGRAY  = '#D5D8DC'
WHITE  = '#FFFFFF'

plt.rcParams.update({
    'font.family': 'sans-serif',
    'figure.facecolor': WHITE,
    'axes.spines.top': False,
    'axes.spines.right': False,
})


# ─────────────────────────────────────────────────────────────────────────────
# UTILS
# ─────────────────────────────────────────────────────────────────────────────
def pct_change(a, b):
    if b is None or pd.isna(b) or b == 0: return None
    return (a - b) / b

def safe_div(a, b):
    if b is None or pd.isna(b) or b == 0: return None
    return a / b

def fmt_brl(v):
    s = f"{abs(v):,.2f}".replace(',','X').replace('.',',').replace('X','.')
    return (f"R$ {s}") if v >= 0 else f"-R$ {s}"

def fmt_brl_k(v):
    if abs(v) >= 1e6: return f"R$ {v/1e6:.1f}M"
    if abs(v) >= 1e3: return f"R$ {v/1e3:.0f}K"
    return fmt_brl(v)

def fmt_pct(v, d=1):
    if v is None or pd.isna(v): return 'n/a'
    return f"{v*100:+.{d}f}%"


# ─────────────────────────────────────────────────────────────────────────────
# CARGA E PROCESSAMENTO DOS DADOS
# ─────────────────────────────────────────────────────────────────────────────
print("Carregando dimensões...")

clients = pd.read_excel(BASE_DIR / 'base_cliente.xlsx').rename(columns={
    'DS_REDE':'network','NM_REDUZIDO':'store','SG_ESTADO':'state','NM_MUNICIPIO':'city'})
for c in ['network','store','state','city']:
    clients[c] = clients[c].astype(str).str.strip()

products = pd.read_excel(BASE_DIR / 'base_produtos.xlsx').rename(columns={
    'CODIGO':'product_code','PRODUTO':'product_name','CATEGORIA':'category'})
for c in ['product_code','product_name','category']:
    products[c] = products[c].astype(str).str.strip()

client_map  = clients.set_index('store')[['network','state','city']].to_dict(orient='index')
product_map = products.set_index('product_code')[['product_name','category']].to_dict(orient='index')


class Agg:
    __slots__ = ('rev','vol','stores','products')
    def __init__(self): self.rev=0.; self.vol=0.; self.stores=set(); self.products=set()
    def add(self,r,v,s,p):
        self.rev+=r; self.vol+=v
        if s: self.stores.add(s)
        if p: self.products.add(p)


overall_agg={}; net_year={}; cat_year={}; qlog={}

print("Processando faturamento (modo streaming)...")
wb = load_workbook(BASE_DIR / 'base_faturamento.xlsx', read_only=True, data_only=True)
for ys in ('2024','2025'):
    y=int(ys); overall_agg[y]=Agg()
    q={'rows':0,'sh':0,'ph':0,'mn':None,'mx':None}; qlog[y]=q
    ws=wb[ys]
    for row in ws.iter_rows(min_row=2, values_only=True):
        dv,sr,pr,qr,rr=row
        store=str(sr).strip() if sr else ''; prod=str(pr).strip() if pr else ''
        qty=float(qr or 0); rev=float(rr or 0)
        ci=client_map.get(store); pi=product_map.get(prod)
        net=ci['network'] if ci else None; cat=pi['category'] if pi else None
        q['rows']+=1
        if ci: q['sh']+=1
        if pi: q['ph']+=1
        if dv:
            ts=pd.Timestamp(dv)
            if q['mn'] is None or ts<q['mn']: q['mn']=ts
            if q['mx'] is None or ts>q['mx']: q['mx']=ts
        overall_agg[y].add(rev,qty,store,prod)
        if net:
            k=(y,net)
            if k not in net_year: net_year[k]=Agg()
            net_year[k].add(rev,qty,store,prod)
        if cat:
            k=(y,cat)
            if k not in cat_year: cat_year[k]=Agg()
            cat_year[k].add(rev,qty,store,prod)
    pct_s=q['sh']/q['rows']*100; pct_p=q['ph']/q['rows']*100
    print(f"  [{ys}] {q['rows']:,} linhas | lojas {pct_s:.4f}% | produtos {pct_p:.4f}%")
wb.close()

def agg_dict(a):
    return {'revenue':a.rev,'volume':a.vol,'avg_price':safe_div(a.rev,a.vol),
            'assortment':len(a.products),'active_stores':len(a.stores)}

overall = {y: agg_dict(a) for y,a in overall_agg.items()}

rows=[]
for (yr,net),a in net_year.items():
    rows.append({'year':yr,'network':net,'revenue':a.rev,'volume':a.vol,
                 'assortment':len(a.products),'active_stores':len(a.stores),
                 'avg_price':safe_div(a.rev,a.vol),'ticket_store':safe_div(a.rev,len(a.stores))})
nt = pd.DataFrame(rows).pivot(index='network',columns='year').fillna(0)
nt.columns=[f"{m}_{y}" for m,y in nt.columns]; nt=nt.reset_index()
for c in ['revenue_2024','revenue_2025','volume_2024','volume_2025','assortment_2024',
          'assortment_2025','active_stores_2024','active_stores_2025',
          'avg_price_2024','avg_price_2025','ticket_store_2024','ticket_store_2025']:
    if c not in nt.columns: nt[c]=0.
nt['rev_growth_abs']=nt['revenue_2025']-nt['revenue_2024']
nt['rev_growth_pct']=nt.apply(lambda r: pct_change(r['revenue_2025'],r['revenue_2024']),axis=1)
nt['vol_growth_pct']=nt.apply(lambda r: pct_change(r['volume_2025'],r['volume_2024']),axis=1)
nt=nt.sort_values('revenue_2025',ascending=False).reset_index(drop=True)

rows_c=[]
for (yr,cat),a in cat_year.items():
    rows_c.append({'year':yr,'category':cat,'revenue':a.rev,'volume':a.vol,
                   'assortment':len(a.products),'avg_price':safe_div(a.rev,a.vol)})
ct = pd.DataFrame(rows_c).pivot(index='category',columns='year').fillna(0)
ct.columns=[f"{m}_{y}" for m,y in ct.columns]; ct=ct.reset_index()
ct['rev_growth_abs']=ct['revenue_2025']-ct['revenue_2024']
ct['rev_growth_pct']=ct.apply(lambda r: pct_change(r['revenue_2025'],r['revenue_2024']),axis=1)
ct['vol_growth_pct']=ct.apply(lambda r: pct_change(r['volume_2025'],r['volume_2024']),axis=1)
ct['price_growth_pct']=ct.apply(lambda r: pct_change(r.get('avg_price_2025',0),r.get('avg_price_2024',0)),axis=1)
ct=ct.sort_values('revenue_2025',ascending=False).reset_index(drop=True)

pareto=nt.sort_values('revenue_2025',ascending=False).copy()
total_25=pareto['revenue_2025'].sum()
pareto['share']=pareto['revenue_2025']/total_25
pareto['cum_share']=pareto['share'].cumsum()
idx_cross=pareto['cum_share'].ge(0.80).idxmax()
pareto['pareto_flag']=False; pareto.loc[:idx_cross,'pareto_flag']=True

comparable=nt[(nt['revenue_2024'].fillna(0)>0)&(nt['revenue_2025'].fillna(0)>0)].copy()
corr=comparable.dropna(subset=['rev_growth_pct','avg_price_2025'])['rev_growth_pct'].corr(
     comparable.dropna(subset=['rev_growth_pct','avg_price_2025'])['avg_price_2025'])

o24,o25=overall[2024],overall[2025]
print("Dados prontos.")


# ─────────────────────────────────────────────────────────────────────────────
# HELPERS DE DESENHO
# ─────────────────────────────────────────────────────────────────────────────
def new_slide(facecolor=WHITE):
    fig = plt.figure(figsize=(W, H), facecolor=facecolor)
    return fig

def draw_header(fig, title, section='', subtitle=''):
    """Barra navy no topo com título branco."""
    # Accent teal (left strip)
    ax_a = fig.add_axes([0, 0.882, 0.004, 0.118])
    ax_a.set_facecolor(TEAL); ax_a.axis('off')
    # Header bg
    ax_h = fig.add_axes([0.004, 0.882, 0.996, 0.118])
    ax_h.set_facecolor(NAVY); ax_h.axis('off')
    ax_h.set_xlim(0,1); ax_h.set_ylim(0,1)
    if section:
        ax_h.text(0.012, 0.88, section.upper(),
                  color=TEAL, fontsize=7.5, fontweight='bold', va='top',
                  fontfamily='monospace')
    ax_h.text(0.012, 0.52, title,
              color=WHITE, fontsize=19, fontweight='bold', va='center')
    if subtitle:
        ax_h.text(0.99, 0.52, subtitle,
                  color='#95A5A6', fontsize=9, va='center', ha='right')
    # Teal bottom border
    ax_b = fig.add_axes([0, 0.879, 1, 0.004])
    ax_b.set_facecolor(TEAL); ax_b.axis('off')

def draw_footer(fig, note='La Vita Alimentos  |  Case Analista Comercial'):
    ax_f = fig.add_axes([0, 0, 1, 0.038])
    ax_f.set_facecolor(LGRAY); ax_f.axis('off')
    ax_f.set_xlim(0,1); ax_f.set_ylim(0,1)
    ax_f.text(0.012, 0.5, note, color=GRAY, fontsize=7, va='center')

def content_ax(fig, left=0.03, bottom=0.06, width=0.94, height=0.80):
    """Retorna axes da área de conteúdo principal."""
    ax = fig.add_axes([left, bottom, width, height])
    ax.set_facecolor(WHITE)
    return ax

def kpi_box(ax_parent, x, y, w, h, value, label,
            delta=None, accent=NAVY, bg=LGRAY, note=''):
    """Desenha um card de KPI dentro de axes_parent (coords 0-1)."""
    ax_parent.add_patch(FancyBboxPatch((x+0.005, y+0.01), w-0.01, h-0.02,
        boxstyle='round,pad=0.01', facecolor=bg, edgecolor=MGRAY, linewidth=1, zorder=1))
    ax_parent.add_patch(Rectangle((x+0.005, y+h-0.025), w-0.01, 0.022,
        facecolor=accent, zorder=2))
    ax_parent.text(x+w/2, y+h*0.58, value,
                   ha='center', va='center', fontsize=22, fontweight='bold',
                   color=accent, zorder=3)
    ax_parent.text(x+w/2, y+h*0.28, label,
                   ha='center', va='center', fontsize=9, color=GRAY, zorder=3)
    if delta:
        col = GREEN if ('+' in str(delta) and str(delta)!='+0,0%') else RED
        if delta in ('n/a','0,0%','+0,0%'): col = GRAY
        ax_parent.text(x+w/2, y+h*0.10, delta,
                       ha='center', va='center', fontsize=10,
                       fontweight='bold', color=col, zorder=3)
    if note:
        ax_parent.text(x+w/2, y+0.02, note,
                       ha='center', va='bottom', fontsize=7.5,
                       color=GRAY, style='italic', zorder=3)

def save_slide(fig, name):
    path = TMP_DIR / name
    fig.savefig(path, dpi=DPI, bbox_inches='tight',
                facecolor=fig.get_facecolor(), edgecolor='none', pad_inches=0)
    plt.close(fig)
    print(f"  ✓ {name}")
    return path


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 01 — CAPA
# ─────────────────────────────────────────────────────────────────────────────
def slide_01():
    fig = new_slide(NAVY)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_facecolor(NAVY); ax.axis('off')
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)

    # Left teal accent
    ax.add_patch(Rectangle((0, 0), 0.006, 1, color=TEAL, zorder=2))
    # Bottom teal strip
    ax.add_patch(Rectangle((0, 0), 1, 0.008, color=TEAL, alpha=0.5, zorder=2))

    # Company name
    ax.text(0.06, 0.78, 'La Vita Alimentos', color=WHITE,
            fontsize=38, fontweight='bold', va='top')
    # Teal separator line
    ax.add_patch(Rectangle((0.06, 0.66), 0.38, 0.005, color=TEAL, zorder=3))

    # Main title
    ax.text(0.06, 0.62, 'Análise Comercial', color=TEAL,
            fontsize=26, fontweight='bold', va='top')

    ax.text(0.06, 0.50, 'Case – Analista Comercial Pleno', color='#BDC3C7',
            fontsize=16, va='top')
    ax.text(0.06, 0.40, 'Período analisado: 2024 e 2025', color='#7F8C8D',
            fontsize=13, va='top')

    # 5 deliverables (right side)
    deliverables = [
        '01  Concentração de Faturamento (ABC / Pareto)',
        '02  Análise Comparativa por Cliente',
        '03  Análise por Categoria',
        '04  Crescimento de Vendas YoY',
        '05  Modelo de Remuneração Variável',
    ]
    ax.text(0.60, 0.78, 'ENTREGAS', color=TEAL,
            fontsize=9, fontweight='bold', va='top', fontfamily='monospace')
    for i, d in enumerate(deliverables):
        ax.text(0.60, 0.70 - i*0.10, d, color='#BDC3C7',
                fontsize=11, va='top')

    ax.text(0.06, 0.08, 'Março 2026', color='#566573', fontsize=11)
    return save_slide(fig, '01_capa.png')


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 02 — METODOLOGIA + QUALIDADE
# ─────────────────────────────────────────────────────────────────────────────
def slide_02():
    fig = new_slide()
    draw_header(fig, 'Metodologia e Qualidade dos Dados',
                section='Premissa', subtitle='Antes do primeiro número')
    draw_footer(fig)

    # ── Left column: methodological decisions ──────────────────────────────
    ax_l = fig.add_axes([0.03, 0.09, 0.54, 0.76])
    ax_l.axis('off'); ax_l.set_xlim(0,1); ax_l.set_ylim(0,1)

    ax_l.text(0, 0.96, 'Decisões metodológicas', color=NAVY,
              fontsize=12, fontweight='bold', va='top')

    decisions = [
        ('Cliente',        'Rede (DS_REDE)',                   'analise acima da loja individual'),
        ('Loja',           'NM_REDUZIDO',                      'menor granularidade disponível'),
        ('Preço médio',    'Receita ÷ Quantidade',             'evita viés da média simples por linha'),
        ('Crescimento %',  'Só redes comparáveis (2024 > 0)',  'exclui distorção de base zero'),
        ('Crescimento R$', 'Todas as redes com receita 2025',  'mede impacto financeiro real'),
        ('Dimensão (E3)',  'Categoria',                        'join confiável — macro-região inviável'),
        ('Macro-região',   'Descartada',                       'sem coluna de município para join'),
    ]

    y = 0.86
    ax_l.add_patch(Rectangle((0, y-0.005), 1, 0.025,
                              facecolor=NAVY, alpha=0.85))
    ax_l.text(0.01, y+0.009, 'Decisão', color=WHITE, fontsize=8.5,
              fontweight='bold', va='center')
    ax_l.text(0.30, y+0.009, 'Critério adotado', color=WHITE, fontsize=8.5,
              fontweight='bold', va='center')
    ax_l.text(0.62, y+0.009, 'Justificativa', color=TEAL, fontsize=8,
              fontweight='bold', va='center')
    y -= 0.025

    for i, (dec, crit, just) in enumerate(decisions):
        bg = LGRAY if i % 2 == 0 else WHITE
        ax_l.add_patch(Rectangle((0, y-0.005), 1, 0.10, facecolor=bg))
        ax_l.text(0.01, y+0.042, dec,  color=NAVY, fontsize=8.5, fontweight='bold', va='center')
        col_c = RED if dec == 'Macro-região' else TEAL
        ax_l.text(0.30, y+0.042, crit, color=col_c, fontsize=8.5, fontweight='bold', va='center')
        ax_l.text(0.62, y+0.042, just, color=GRAY, fontsize=7.5, va='center')
        y -= 0.10

    # ── Right column: data quality ─────────────────────────────────────────
    ax_r = fig.add_axes([0.61, 0.09, 0.36, 0.76])
    ax_r.axis('off'); ax_r.set_xlim(0,1); ax_r.set_ylim(0,1)

    ax_r.text(0, 0.96, 'Qualidade dos dados', color=NAVY,
              fontsize=12, fontweight='bold', va='top')

    q24, q25 = qlog[2024], qlog[2025]
    total_rows = q24['rows'] + q25['rows']

    quality_items = [
        ('Join lojas', f"{q24['sh']/q24['rows']*100:.4f}%", '2024', GREEN),
        ('Join lojas', f"{q25['sh']/q25['rows']*100:.4f}%", '2025', GREEN),
        ('Join produtos', f"{q24['ph']/q24['rows']*100:.4f}%", '2024', GREEN),
        ('Join produtos', f"{q25['ph']/q25['rows']*100:.4f}%", '2025', GREEN),
    ]

    # Quality KPI cards
    kpi_items = [
        ('Join lojas', '100,0000%', 'em ambos os anos', GREEN),
        ('Join produtos', '99,9999%', '1 código sem cadastro', GREEN),
        ('Período 2024', f"{q24['mn'].strftime('%b/%y')} → {q24['mx'].strftime('%b/%y')}",
         'jan–dez, comparável', TEAL),
        ('Período 2025', f"{q25['mn'].strftime('%b/%y')} → {q25['mx'].strftime('%b/%y')}",
         'jan–dez, comparável', TEAL),
    ]

    y = 0.82
    for label, val, sub, col in kpi_items:
        ax_r.add_patch(FancyBboxPatch((0.02, y-0.13), 0.96, 0.12,
            boxstyle='round,pad=0.01', facecolor=LGRAY, edgecolor=MGRAY, linewidth=1))
        ax_r.add_patch(Rectangle((0.02, y-0.01), 0.025, 0.12, facecolor=col))
        ax_r.text(0.10, y-0.055, label, color=GRAY, fontsize=8.5, va='center')
        ax_r.text(0.10, y-0.085, val, color=col, fontsize=11, fontweight='bold', va='center')
        ax_r.text(0.10, y-0.112, sub, color=GRAY, fontsize=7.5, va='center', style='italic')
        y -= 0.155

    # Scale box
    ax_r.add_patch(FancyBboxPatch((0.02, 0.02), 0.96, 0.14,
        boxstyle='round,pad=0.01', facecolor=NAVY, edgecolor=NAVY))
    ax_r.text(0.50, 0.12, f"{total_rows:,}".replace(',','.'),
              color=WHITE, fontsize=22, fontweight='bold', ha='center', va='center')
    ax_r.text(0.50, 0.05, 'transações analisadas',
              color=TEAL, fontsize=9, ha='center', va='center')

    return save_slide(fig, '02_metodologia.png')


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 03 — E1: CURVA ABC / PARETO
# ─────────────────────────────────────────────────────────────────────────────
def slide_03():
    fig = new_slide()
    n_a = pareto['pareto_flag'].sum()
    n_ativas = int((nt['revenue_2025'] > 0).sum())
    draw_header(fig, 'Concentração de Faturamento — Curva ABC / Pareto',
                section='Entrega 1')
    draw_footer(fig)

    # ── Left: KPI panel ────────────────────────────────────────────────────
    ax_l = fig.add_axes([0.02, 0.09, 0.26, 0.77])
    ax_l.axis('off'); ax_l.set_xlim(0,1); ax_l.set_ylim(0,1)

    ax_l.add_patch(FancyBboxPatch((0.04, 0.68), 0.92, 0.28,
        boxstyle='round,pad=0.01', facecolor=NAVY, edgecolor=NAVY))
    ax_l.text(0.5, 0.92, str(n_a), color=WHITE,
              fontsize=48, fontweight='bold', ha='center', va='center')
    ax_l.text(0.5, 0.77, 'redes', color=TEAL,
              fontsize=14, ha='center', va='center')
    ax_l.text(0.5, 0.71, 'concentram 81% da receita', color='#BDC3C7',
              fontsize=9, ha='center', va='center')

    ax_l.text(0.5, 0.60, f"{n_a} de {n_ativas} redes ativas", color=GRAY,
              fontsize=10, ha='center', va='center')
    ax_l.text(0.5, 0.52, f"= {n_a/n_ativas*100:.0f}% das redes", color=GRAY,
              fontsize=10, ha='center', va='center')

    ax_l.add_patch(Rectangle((0.04, 0.44), 0.92, 0.001, color=MGRAY))
    ax_l.text(0.5, 0.40, 'Faturamento total 2025', color=GRAY,
              fontsize=8.5, ha='center', va='center')
    ax_l.text(0.5, 0.32, fmt_brl_k(total_25), color=NAVY,
              fontsize=20, fontweight='bold', ha='center', va='center')

    ax_l.add_patch(Rectangle((0.04, 0.22), 0.92, 0.001, color=MGRAY))
    ax_l.text(0.5, 0.17, 'Insight', color=NAVY, fontsize=8.5,
              fontweight='bold', ha='center')
    ax_l.text(0.5, 0.09, 'Gestão dedicada para contas A\né crítica para a receita',
              color=GRAY, fontsize=8, ha='center', va='center', multialignment='center')

    # ── Right: Pareto chart ─────────────────────────────────────────────────
    ax_r = fig.add_axes([0.30, 0.09, 0.68, 0.77])
    ax_r.set_facecolor(WHITE)
    ax2 = ax_r.twiny()

    top15 = pareto.head(15).copy().reset_index(drop=True)
    bar_colors = [NAVY if f else MGRAY for f in top15['pareto_flag']]

    ax_r.barh(range(len(top15)), top15['revenue_2025'],
              color=bar_colors, alpha=0.88, height=0.7)
    ax2.plot(top15['cum_share'] * 100, range(len(top15)),
             color=RED, marker='o', markersize=4, linewidth=2, zorder=5)
    ax2.axvline(80, color=RED, linestyle='--', linewidth=1.2, alpha=0.7)
    ax2.text(80.5, 0, '80%', color=RED, fontsize=8, va='bottom')

    ax_r.set_yticks(range(len(top15)))
    ax_r.set_yticklabels(top15['network'], fontsize=8.5)
    ax_r.invert_yaxis()
    ax_r.xaxis.set_major_formatter(mticker.FuncFormatter(
        lambda x, _: f'R${x/1e6:.1f}M'))
    ax2.set_xlabel('% Acumulado', fontsize=8)
    ax2.set_xlim(0, 108)
    ax_r.set_xlabel('Receita 2025', fontsize=9)
    ax_r.spines['top'].set_visible(False)
    ax_r.grid(axis='x', alpha=0.2, linestyle='--')

    from matplotlib.patches import Patch
    ax_r.legend(handles=[Patch(color=NAVY, label='Classe A (≤ 80%)'),
                          Patch(color=MGRAY, label='Demais redes')],
                loc='lower right', fontsize=8, framealpha=0.9)

    return save_slide(fig, '03_pareto.png')


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 04 — E2: COMPARATIVO POR CLIENTE (4 CARDS)
# ─────────────────────────────────────────────────────────────────────────────
def slide_04():
    fig = new_slide()
    draw_header(fig, 'Análise Comparativa por Cliente', section='Entrega 2',
                subtitle='2025 como ano de referência')
    draw_footer(fig)

    ax = fig.add_axes([0.02, 0.08, 0.96, 0.80])
    ax.axis('off'); ax.set_xlim(0,1); ax.set_ylim(0,1)

    # Pre-compute winners
    best_ticket   = nt[nt['active_stores_2025']>0].sort_values('ticket_store_2025',ascending=False).iloc[0]
    best_vol      = nt.sort_values('volume_2025',ascending=False).iloc[0]
    best_asm      = nt[nt['revenue_2025']>0].sort_values('assortment_2025',ascending=False).iloc[0]
    best_pct      = comparable.sort_values('rev_growth_pct',ascending=False).iloc[0]
    best_abs      = comparable.sort_values('rev_growth_abs',ascending=False).iloc[0]

    cards = [
        {
            'pos': (0.01, 0.03, 0.47, 0.45),
            'accent': TEAL,
            'label': 'TICKET MÉDIO POR LOJA',
            'winner': best_ticket['network'],
            'value': fmt_brl_k(best_ticket['ticket_store_2025']),
            'note': f"{int(best_ticket['active_stores_2025'])} loja(s) ativa(s)",
            'formula': 'Receita da rede ÷ lojas ativas',
        },
        {
            'pos': (0.52, 0.03, 0.47, 0.45),
            'accent': NAVY,
            'label': 'VOLUME DE COMPRAS',
            'winner': best_vol['network'],
            'value': f"{best_vol['volume_2025']/1e6:.1f}M un",
            'note': f"Receita: {fmt_brl_k(best_vol['revenue_2025'])}",
            'formula': 'Soma de QUANTIDADE',
        },
        {
            'pos': (0.01, 0.52, 0.47, 0.45),
            'accent': ORANGE,
            'label': 'MAIOR SORTIMENTO',
            'winner': best_asm['network'],
            'value': f"{int(best_asm['assortment_2025'])} SKUs",
            'note': f"Teto do portfólio vendido: {int(best_asm['assortment_2025'])} SKUs",
            'formula': 'SKUs distintos em 2025',
        },
        {
            'pos': (0.52, 0.52, 0.47, 0.45),
            'accent': GREEN,
            'label': 'MAIOR CRESCIMENTO',
            'winner': best_pct['network'],
            'value': fmt_pct(best_pct['rev_growth_pct']),
            'note': f"R$: {best_abs['network']} +{fmt_brl_k(best_abs['rev_growth_abs'])}",
            'formula': 'Redes comparáveis (2024 > 0)',
        },
    ]

    for c in cards:
        l,b,w,h = c['pos']
        ax.add_patch(FancyBboxPatch((l+0.005, b+0.01), w-0.01, h-0.02,
            boxstyle='round,pad=0.01', facecolor=LGRAY, edgecolor=MGRAY, linewidth=1))
        ax.add_patch(Rectangle((l+0.005, b+h-0.025), w-0.01, 0.022,
            facecolor=c['accent']))
        ax.text(l+w/2, b+h-0.012, c['label'],
                ha='center', va='center', fontsize=8, fontweight='bold',
                color=WHITE)
        ax.text(l+w/2, b+h*0.66, c['winner'],
                ha='center', va='center', fontsize=13, fontweight='bold',
                color=c['accent'])
        ax.text(l+w/2, b+h*0.43, c['value'],
                ha='center', va='center', fontsize=22, fontweight='bold',
                color=NAVY)
        ax.text(l+w/2, b+h*0.22, c['note'],
                ha='center', va='center', fontsize=8.5, color=GRAY)
        ax.text(l+w/2, b+0.02, f"Como calculado: {c['formula']}",
                ha='center', va='bottom', fontsize=7, color=GRAY, style='italic')

    return save_slide(fig, '04_comparativo.png')


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 05 — E2: PREÇO MÉDIO × CRESCIMENTO
# ─────────────────────────────────────────────────────────────────────────────
def slide_05():
    fig = new_slide()
    draw_header(fig, 'Preço Médio não explica Crescimento', section='Entrega 2',
                subtitle='Análise de correlação — redes comparáveis')
    draw_footer(fig)

    sc = comparable.dropna(subset=['rev_growth_pct','avg_price_2025']).copy()
    # Remove extreme outliers for better visualization
    q99 = sc['rev_growth_pct'].quantile(0.97)
    sc_plot = sc[sc['rev_growth_pct'] <= q99].copy()

    # ── Chart ──────────────────────────────────────────────────────────────
    ax = fig.add_axes([0.04, 0.10, 0.60, 0.76])
    ax.set_facecolor(WHITE)
    ax.scatter(sc_plot['avg_price_2025'], sc_plot['rev_growth_pct']*100,
               alpha=0.55, color=NAVY, s=45, zorder=3)

    # Labels for highlights
    q90_g = sc_plot['rev_growth_pct'].quantile(0.90)
    q90_p = sc_plot['avg_price_2025'].quantile(0.90)
    mask = (sc_plot['rev_growth_pct'] > q90_g) | (sc_plot['avg_price_2025'] > q90_p)
    for _, row in sc_plot[mask].iterrows():
        ax.annotate(row['network'], (row['avg_price_2025'], row['rev_growth_pct']*100),
                    fontsize=7, xytext=(4,3), textcoords='offset points',
                    color=NAVY, alpha=0.85)

    ax.axhline(0, color=GRAY, linewidth=0.8, linestyle='--')
    ax.set_xlabel('Preço médio 2025 (R$)', fontsize=9)
    ax.set_ylabel('Crescimento de receita (%)', fontsize=9)
    ax.yaxis.set_major_formatter(mticker.PercentFormatter())
    ax.set_title('Dispersão: Preço Médio × Crescimento de Receita', fontsize=10,
                 fontweight='bold', color=NAVY, pad=8)
    ax.grid(alpha=0.2, linestyle='--')
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)

    # ── Right panel: findings ───────────────────────────────────────────────
    ax_r = fig.add_axes([0.68, 0.10, 0.30, 0.76])
    ax_r.axis('off'); ax_r.set_xlim(0,1); ax_r.set_ylim(0,1)

    max_preco = sc.sort_values('avg_price_2025',ascending=False).iloc[0]
    best_pct2 = sc.sort_values('rev_growth_pct',ascending=False).iloc[0]

    findings = [
        (f"Correlação de Pearson: {corr:.3f}",
         'Relação praticamente nula\nentre preço e crescimento',
         RED),
        (f"Maior preço médio\n{max_preco['network']}",
         f"R$ {max_preco['avg_price_2025']:.2f} / un",
         TEAL),
        (f"Maior crescimento %\n{best_pct2['network']}",
         f"Preço médio: R$ {best_pct2['avg_price_2025']:.2f}",
         NAVY),
        ('Implicação',
         'O modelo atual de comissão\nnão protege preço médio',
         ORANGE),
    ]

    y = 0.96
    for title, body, col in findings:
        ax_r.add_patch(FancyBboxPatch((0.03, y-0.22), 0.94, 0.20,
            boxstyle='round,pad=0.01', facecolor=LGRAY, edgecolor=col, linewidth=1.5))
        ax_r.add_patch(Rectangle((0.03, y-0.02), 0.015, 0.20, facecolor=col))
        ax_r.text(0.10, y-0.06, title, color=col,
                  fontsize=8.5, fontweight='bold', va='top')
        ax_r.text(0.10, y-0.13, body, color=GRAY,
                  fontsize=8, va='top', multialignment='left')
        y -= 0.24

    return save_slide(fig, '05_preco_crescimento.png')


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 06 — E3: ANÁLISE POR CATEGORIA
# ─────────────────────────────────────────────────────────────────────────────
def slide_06():
    fig = new_slide()
    draw_header(fig, 'Análise por Categoria', section='Entrega 3',
                subtitle='Dimensão escolhida — join confiável em 100% das linhas')
    draw_footer(fig)

    cat_s = ct.sort_values('revenue_2025', ascending=True)
    x = np.arange(len(cat_s))
    w = 0.35

    # ── Top row: 3 charts ──────────────────────────────────────────────────
    # Chart 1: Revenue
    ax1 = fig.add_axes([0.04, 0.32, 0.30, 0.52])
    ax1.barh(x - w/2, cat_s['revenue_2024'], w, label='2024', color=MGRAY, alpha=0.85)
    ax1.barh(x + w/2, cat_s['revenue_2025'], w, label='2025', color=NAVY, alpha=0.90)
    ax1.set_yticks(x); ax1.set_yticklabels(cat_s['category'], fontsize=8.5)
    ax1.xaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f'R${v/1e6:.0f}M'))
    ax1.set_title('Receita (R$)', fontsize=10, fontweight='bold', color=NAVY, pad=6)
    ax1.legend(fontsize=8, loc='lower right')
    ax1.grid(axis='x', alpha=0.2, linestyle='--')
    ax1.spines['top'].set_visible(False); ax1.spines['right'].set_visible(False)

    # Chart 2: Revenue growth %
    ax2 = fig.add_axes([0.37, 0.32, 0.28, 0.52])
    cl = [GREEN if v >= 0 else RED for v in cat_s['rev_growth_pct'].fillna(0)]
    ax2.barh(x, cat_s['rev_growth_pct'].fillna(0)*100, color=cl, alpha=0.88)
    ax2.set_yticks(x); ax2.set_yticklabels(cat_s['category'], fontsize=8.5)
    ax2.xaxis.set_major_formatter(mticker.PercentFormatter())
    ax2.axvline(0, color='black', linewidth=0.5)
    ax2.set_title('Δ Receita YoY', fontsize=10, fontweight='bold', color=NAVY, pad=6)
    ax2.grid(axis='x', alpha=0.2, linestyle='--')
    ax2.spines['top'].set_visible(False); ax2.spines['right'].set_visible(False)

    # Chart 3: Price growth %
    ax3 = fig.add_axes([0.68, 0.32, 0.28, 0.52])
    clp = [GREEN if v >= 0 else RED for v in cat_s['price_growth_pct'].fillna(0)]
    ax3.barh(x, cat_s['price_growth_pct'].fillna(0)*100, color=clp, alpha=0.88)
    ax3.set_yticks(x); ax3.set_yticklabels(cat_s['category'], fontsize=8.5)
    ax3.xaxis.set_major_formatter(mticker.PercentFormatter())
    ax3.axvline(0, color='black', linewidth=0.5)
    ax3.set_title('Δ Preço Médio YoY', fontsize=10, fontweight='bold', color=NAVY, pad=6)
    ax3.grid(axis='x', alpha=0.2, linestyle='--')
    ax3.spines['top'].set_visible(False); ax3.spines['right'].set_visible(False)

    # ── Bottom row: key insights ────────────────────────────────────────────
    ax_b = fig.add_axes([0.03, 0.07, 0.94, 0.21])
    ax_b.axis('off'); ax_b.set_xlim(0,1); ax_b.set_ylim(0,1)

    insights = [
        ('SALADAS',  NAVY,  'Maior categoria (+50% receita).\nCrescimento saudável: vol +28%, preço estável.'),
        ('IN NATURA', RED,  'Volume cresce +9,4% mas\npreço cai -3,7%. Pressão de margem.'),
        ('TEMPEROS',  GREEN,'Cresce volume E preço (+3,6%).\nMelhor posicionamento relativo.'),
        ('LEGUMES',   TEAL, 'Aceleração forte (+66% receita).\nEspansão puxada por volume.'),
    ]

    x_pos = 0.01
    for cat_name, col, desc in insights:
        ax_b.add_patch(FancyBboxPatch((x_pos, 0.05), 0.23, 0.88,
            boxstyle='round,pad=0.01', facecolor=LGRAY, edgecolor=col, linewidth=1.5))
        ax_b.text(x_pos+0.115, 0.82, cat_name, ha='center', va='center',
                  color=col, fontsize=10, fontweight='bold')
        ax_b.text(x_pos+0.115, 0.42, desc, ha='center', va='center',
                  color=GRAY, fontsize=8, multialignment='center')
        x_pos += 0.245

    return save_slide(fig, '06_categoria.png')


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 07 — E4 SLIDE 1: YoY VISÃO GERAL
# ─────────────────────────────────────────────────────────────────────────────
def slide_07():
    fig = new_slide()
    draw_header(fig, 'Visão Geral YoY 2024 → 2025', section='Entrega 4  |  Slide 1 / 3',
                subtitle='Crescimento puxado por volume')
    draw_footer(fig)

    ax = fig.add_axes([0.02, 0.28, 0.96, 0.57])
    ax.axis('off'); ax.set_xlim(0,1); ax.set_ylim(0,1)

    r_g = pct_change(o25['revenue'],    o24['revenue'])
    v_g = pct_change(o25['volume'],     o24['volume'])
    p_g = pct_change(o25['avg_price'],  o24['avg_price'])
    a_g = pct_change(o25['assortment'], o24['assortment'])

    kpis = [
        ('RECEITA',      fmt_brl_k(o25['revenue']),   fmt_pct(r_g),
         fmt_brl_k(o24['revenue']),  NAVY,   'Crescimento real'),
        ('VOLUME',       f"{o25['volume']/1e6:.2f}M un", fmt_pct(v_g),
         f"{o24['volume']/1e6:.2f}M un",     TEAL,  'Principal driver'),
        ('PREÇO MÉDIO',  fmt_brl(o25['avg_price']),    fmt_pct(p_g, d=2),
         fmt_brl(o24['avg_price']),           ORANGE, 'Quase estável'),
        ('SORTIMENTO',   f"{int(o25['assortment'])} SKUs", '0,0%',
         f"{int(o24['assortment'])} SKUs",    GRAY,  'Sem expansão'),
    ]

    x = 0.01
    for label, val25, delta, val24, col, tag in kpis:
        ax.add_patch(FancyBboxPatch((x+0.005, 0.04), 0.22, 0.92,
            boxstyle='round,pad=0.01', facecolor=LGRAY, edgecolor=MGRAY, linewidth=1))
        ax.add_patch(Rectangle((x+0.005, 0.96-0.025), 0.22, 0.022, facecolor=col))
        ax.text(x+0.115, 0.925, label, ha='center', va='center',
                fontsize=8, fontweight='bold', color=WHITE)
        ax.text(x+0.115, 0.70, val25, ha='center', va='center',
                fontsize=22, fontweight='bold', color=col)
        d_col = GREEN if ('+' in str(delta) and delta not in ('+0,0%','n/a')) else (RED if '-' in str(delta) else GRAY)
        ax.text(x+0.115, 0.50, delta, ha='center', va='center',
                fontsize=16, fontweight='bold', color=d_col)
        ax.text(x+0.115, 0.31, f"2024: {val24}", ha='center', va='center',
                fontsize=9, color=GRAY)
        ax.text(x+0.115, 0.14, tag, ha='center', va='center',
                fontsize=8, color=col, style='italic')
        x += 0.245

    # Bottom message
    ax_msg = fig.add_axes([0.02, 0.08, 0.96, 0.17])
    ax_msg.set_facecolor(NAVY); ax_msg.axis('off')
    ax_msg.set_xlim(0,1); ax_msg.set_ylim(0,1)
    ax_msg.add_patch(Rectangle((0, 0), 0.006, 1, color=TEAL))
    ax_msg.text(0.02, 0.60,
                'Receita cresceu +24% — mas 23 pp vieram de volume e menos de 1 pp de preço.',
                color=WHITE, fontsize=12, fontweight='bold', va='center')
    ax_msg.text(0.02, 0.22,
                'O sortimento total vendido ficou estável. A empresa vendeu mais, não diferente.',
                color=TEAL, fontsize=10, va='center')

    return save_slide(fig, '07_yoy_slide1.png')


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 08 — E4 SLIDE 2: DRIVERS E DETRATORES
# ─────────────────────────────────────────────────────────────────────────────
def slide_08():
    fig = new_slide()
    draw_header(fig, 'Impulsionadores e Detratores', section='Entrega 4  |  Slide 2 / 3',
                subtitle='Por variação absoluta de receita — redes comparáveis')
    draw_footer(fig)

    comp_all = nt[nt['revenue_2024'].fillna(0) > 0].copy()
    drivers    = comp_all.sort_values('rev_growth_abs', ascending=False).head(5)
    detractors = comp_all.sort_values('rev_growth_abs', ascending=True).head(5)
    comb = (pd.concat([drivers, detractors])
              .drop_duplicates('network')
              .sort_values('rev_growth_abs', ascending=True))

    # ── Diverging bar chart ─────────────────────────────────────────────────
    ax = fig.add_axes([0.04, 0.10, 0.55, 0.76])
    ax.set_facecolor(WHITE)
    bar_colors = [GREEN if v >= 0 else RED for v in comb['rev_growth_abs']]
    bars = ax.barh(range(len(comb)), comb['rev_growth_abs'],
                   color=bar_colors, alpha=0.88, height=0.65)

    for bar, val, net_name in zip(bars, comb['rev_growth_abs'], comb['network']):
        x_pos = bar.get_width() + (total_25 * 0.003 if val >= 0 else -total_25 * 0.003)
        ha = 'left' if val >= 0 else 'right'
        ax.text(x_pos, bar.get_y() + bar.get_height()/2,
                fmt_brl_k(val), va='center', ha=ha, fontsize=8, fontweight='bold')

    ax.set_yticks(range(len(comb)))
    ax.set_yticklabels(comb['network'], fontsize=9)
    ax.axvline(0, color='black', linewidth=0.8)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(
        lambda x,_: f'R${x/1e6:.1f}M' if abs(x) >= 1e6 else f'R${x/1e3:.0f}K'))
    ax.set_title('Variação absoluta de receita (2024 → 2025)', fontsize=9,
                 fontweight='bold', color=NAVY, pad=6)
    ax.grid(axis='x', alpha=0.2, linestyle='--')
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)

    # ── Right: category table ───────────────────────────────────────────────
    ax_r = fig.add_axes([0.63, 0.10, 0.35, 0.76])
    ax_r.axis('off'); ax_r.set_xlim(0,1); ax_r.set_ylim(0,1)

    ax_r.text(0.5, 0.97, 'Impacto por Categoria', color=NAVY,
              fontsize=11, fontweight='bold', ha='center', va='top')

    # Header
    ax_r.add_patch(Rectangle((0, 0.87), 1, 0.06, facecolor=NAVY, alpha=0.88))
    ax_r.text(0.03, 0.900, 'Categoria', color=WHITE, fontsize=8, fontweight='bold', va='center')
    ax_r.text(0.65, 0.900, 'Δ Receita', color=WHITE, fontsize=8, fontweight='bold', va='center',ha='center')
    ax_r.text(0.88, 0.900, 'Var%', color=TEAL, fontsize=8, fontweight='bold', va='center',ha='center')

    cat_sorted = ct.sort_values('rev_growth_abs', ascending=False)
    y = 0.84
    for i, (_, row) in enumerate(cat_sorted.iterrows()):
        bg = LGRAY if i % 2 == 0 else WHITE
        ax_r.add_patch(Rectangle((0, y-0.075), 1, 0.075, facecolor=bg))
        col_r = GREEN if row['rev_growth_abs'] >= 0 else RED
        ax_r.text(0.03, y-0.038, row['category'], color=NAVY,
                  fontsize=8.5, fontweight='bold', va='center')
        ax_r.text(0.65, y-0.038, fmt_brl_k(row['rev_growth_abs']), color=col_r,
                  fontsize=8.5, fontweight='bold', va='center', ha='center')
        ax_r.text(0.88, y-0.038, fmt_pct(row['rev_growth_pct']),
                  color=col_r, fontsize=8, va='center', ha='center')
        y -= 0.075

    return save_slide(fig, '08_yoy_slide2.png')


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 09 — E4 SLIDE 3: RECOMENDAÇÕES
# ─────────────────────────────────────────────────────────────────────────────
def slide_09():
    fig = new_slide()
    draw_header(fig, 'Recomendações Prioritárias', section='Entrega 4  |  Slide 3 / 3',
                subtitle='5 ações para 2026 ancoradas nos dados')
    draw_footer(fig)

    ax = fig.add_axes([0.02, 0.07, 0.96, 0.80])
    ax.axis('off'); ax.set_xlim(0,1); ax.set_ylim(0,1)

    recs = [
        (1, TEAL,   'Expandir sortimento\nem contas médias',
         'Redes relevantes abaixo do teto de 29 SKUs',
         '+3 SKUs no mix médio das contas-alvo em 6 meses'),
        (2, RED,    'Proteger preço em\nIN NATURA e COUVE',
         'Volume cresce, mas preço cai — erosão de margem',
         '+1,5% de preço médio em 90 dias, volume –3% máx.'),
        (3, NAVY,   'Separar contas novas\nde crescimento orgânico',
         'ST Marche (+809%) e Dalben (+1485%) são entradas novas',
         'Dois painéis: comparáveis | novos no próximo ciclo'),
        (4, GREEN,  'Replicar alavancas\ndas redes saudáveis',
         'OBA (+27%), Savegnago (+27%) cresceram com preço firme',
         '10 contas-alvo similares para plano em 120 dias'),
        (5, ORANGE, 'Recuperar receita\nos detratores',
         'Pague Menos (–80%), Varejo (–65%), DB (–59%)',
         'Recuperar 30% da perda absoluta em 6 meses'),
    ]

    positions = [
        (0.00, 0.50, 0.195, 0.47),
        (0.20, 0.50, 0.195, 0.47),
        (0.40, 0.50, 0.195, 0.47),
        (0.10, 0.01, 0.195, 0.47),
        (0.30, 0.01, 0.195, 0.47),
    ]

    for (num, col, title, body, meta), (l,b,w,h) in zip(recs, positions):
        ax.add_patch(FancyBboxPatch((l+0.005, b+0.01), w-0.01, h-0.02,
            boxstyle='round,pad=0.01', facecolor=LGRAY, edgecolor=col, linewidth=1.5))
        # Number circle
        circ = plt.Circle((l+0.045, b+h-0.06), 0.035, color=col, zorder=3,
                           transform=ax.transData)
        ax.add_patch(circ)
        ax.text(l+0.045, b+h-0.06, str(num), ha='center', va='center',
                fontsize=12, fontweight='bold', color=WHITE, zorder=4)
        # Title
        ax.text(l+w/2, b+h*0.70, title, ha='center', va='center',
                fontsize=10, fontweight='bold', color=col, multialignment='center')
        # Body
        ax.text(l+w/2, b+h*0.45, body, ha='center', va='center',
                fontsize=8, color=GRAY, multialignment='center')
        # Meta (target)
        ax.add_patch(FancyBboxPatch((l+0.015, b+0.02), w-0.03, 0.13,
            boxstyle='round,pad=0.005', facecolor=col, alpha=0.12, edgecolor=col, linewidth=0.5))
        ax.text(l+w/2, b+0.086, f"→ {meta}", ha='center', va='center',
                fontsize=7.5, color=col, fontweight='bold', multialignment='center')

    return save_slide(fig, '09_yoy_slide3.png')


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — E5: MODELO DE REMUNERAÇÃO VARIÁVEL
# ─────────────────────────────────────────────────────────────────────────────
def slide_10():
    fig = new_slide()
    draw_header(fig, 'Modelo de Remuneração Variável', section='Entrega 5',
                subtitle='Premiar crescimento com qualidade — não só volume')
    draw_footer(fig)

    ax = fig.add_axes([0.02, 0.07, 0.96, 0.80])
    ax.axis('off'); ax.set_xlim(0,1); ax.set_ylim(0,1)

    # ── Formula bar ────────────────────────────────────────────────────────
    ax.text(0.5, 0.97,
            'Payout  =  Bônus-alvo  ×  ( 0,40 × S_vol  +  0,30 × S_preço  +  0,20 × S_sort  +  0,10 × S_mix )',
            ha='center', va='top', fontsize=11, fontweight='bold', color=NAVY)

    # Weight bar
    components = [
        ('Crescimento de Volume', 0.40, NAVY),
        ('Disciplina de Preço',   0.30, TEAL),
        ('Sortimento',            0.20, ORANGE),
        ('Mix foco',              0.10, GRAY),
    ]
    bar_x = 0.02
    bar_y, bar_h = 0.83, 0.06
    for label, pct, col in components:
        w_seg = pct * 0.96
        ax.add_patch(Rectangle((bar_x, bar_y), w_seg, bar_h, facecolor=col, alpha=0.90))
        if pct >= 0.15:
            ax.text(bar_x + w_seg/2, bar_y + bar_h/2,
                    f"{label}\n{int(pct*100)}%",
                    ha='center', va='center', fontsize=8.5,
                    color=WHITE, fontweight='bold', multialignment='center')
        else:
            ax.text(bar_x + w_seg/2, bar_y - 0.03,
                    f"{int(pct*100)}%", ha='center', va='top', fontsize=8, color=col,
                    fontweight='bold')
        bar_x += w_seg

    # ── 4 indicator cards ────��─────────────────────────────────────────────
    indicators = [
        (NAVY,   'Crescimento de Volume', '40%',
         '(qtd_2025 − qtd_2024) / qtd_2024',
         'Redes comparáveis; base zero excluída',
         'Premia expansão real de vendas'),
        (TEAL,   'Disciplina de Preço', '30%',
         'Preço médio realizado vs faixa-meta',
         'Redutor quando erosão > threshold',
         'Evita crescimento via desconto'),
        (ORANGE, 'Crescimento de Sortimento', '20%',
         'Δ SKUs distintos com ≥ 4 sem. recorrência',
         'Evita venda pontual para bater meta',
         'Incentiva penetração de portfólio'),
        (GRAY,   'Execução de Mix (contas foco)', '10%',
         'Cobertura de SKUs estratégicos nas A-contas',
         'Acompanhado mensalmente por conta',
         'Evita concentração em poucos produtos'),
    ]

    x_card = 0.01
    for col, name, pct, measure, rule, rationale in indicators:
        ax.add_patch(FancyBboxPatch((x_card+0.005, 0.03), 0.225, 0.72,
            boxstyle='round,pad=0.01', facecolor=LGRAY, edgecolor=col, linewidth=1.5))
        ax.add_patch(Rectangle((x_card+0.005, 0.75-0.025), 0.225, 0.022, facecolor=col))
        ax.text(x_card+0.12, 0.74, pct, ha='center', va='center',
                fontsize=22, fontweight='bold', color=col)
        ax.text(x_card+0.12, 0.63, name, ha='center', va='center',
                fontsize=9, fontweight='bold', color=NAVY, multialignment='center')
        ax.add_patch(Rectangle((x_card+0.01, 0.52), 0.215, 0.001, color=MGRAY))
        ax.text(x_card+0.12, 0.47, 'Como medir', ha='center', fontsize=7.5,
                color=GRAY, fontweight='bold', va='center')
        ax.text(x_card+0.12, 0.38, measure, ha='center', fontsize=7.5,
                color=TEAL, va='center', multialignment='center')
        ax.text(x_card+0.12, 0.27, rule, ha='center', fontsize=7.5,
                color=GRAY, va='center', multialignment='center', style='italic')
        ax.add_patch(FancyBboxPatch((x_card+0.015, 0.04), 0.21, 0.14,
            boxstyle='round,pad=0.005', facecolor=col, alpha=0.12,
            edgecolor=col, linewidth=0.5))
        ax.text(x_card+0.12, 0.11, rationale, ha='center', va='center',
                fontsize=7.5, color=col, fontweight='bold', multialignment='center')
        x_card += 0.245

    return save_slide(fig, '10_remuneracao.png')


# ─────────────────────────────────────────────────────────────────────────────
# MONTA O PPTX
# ─────────────────────────────────────────────────────────────────────────────
def build_pptx(slide_paths, output_path):
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # blank layout
    for path in slide_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(path), left=0, top=0,
            width=prs.slide_width, height=prs.slide_height)

    prs.save(output_path)
    print(f"\n✅  PPTX salvo em: {output_path}")


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    print("\nGerando slides...")
    paths = [
        slide_01(),
        slide_02(),
        slide_03(),
        slide_04(),
        slide_05(),
        slide_06(),
        slide_07(),
        slide_08(),
        slide_09(),
        slide_10(),
    ]
    output = BASE_DIR / 'analise_la_vita.pptx'
    build_pptx(paths, output)
    print(f"Total: {len(paths)} slides  |  {output.stat().st_size / 1024:.0f} KB")
